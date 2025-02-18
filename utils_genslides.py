import json
from pptx import Presentation
import random
import datetime
from pptx.util import Inches
from PIL import Image
import jieba

# TO-DO:  These modules need refactoring.  They are a bit messy with all the offsets.
# This is due to the placeholder index values being hardcoded in the powerpoint templates.
# The other complexity is that even for placeholders we don't use, we have to loop
# through them to set them to empty strings.  This is because the placeholders already
# have text in them from the template.  


# Load json file with the lesson content into a dictionary
def load_json(json_file):
    with open(json_file) as json_file:
        lesson_content = json.load(json_file)
    return lesson_content

# Add Vocabulary Summary Page
def addslide_title_teacher(prs, data):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.placeholders[13].text = data.get("title")
    slide.placeholders[14].text = data.get("subtitle")

def addslide_title_student(prs, data):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.placeholders[13].text = data.get("title")
    slide.placeholders[14].text = data.get("subtitle")

def addslide_comicscript(prs, data):
    slide_layout = prs.slide_layouts[2]
    offsetA = 13
    offsetB = 20
    offsetC = 34
    offsetD = 48

    for i, eachline in enumerate(data["script_list"]):
        if i%7 == 0:
            slide = prs.slides.add_slide(slide_layout)
            k = 0
        slide.placeholders[offsetA + k].text = eachline.get("script_text")

        for j in [0,1]:
            try:    
                each_vocab = eachline.get("vocab")[j]
                slide.placeholders[offsetB + 2*k +j].text = each_vocab.get("SC_vocab")
                slide.placeholders[offsetC + 2*k +j].text = each_vocab.get("TC_vocab")
                slide.placeholders[offsetD + 2*k +j].text = each_vocab.get("english_meaning")
            except:
                slide.placeholders[offsetB + 2*k +j].text = " "
                slide.placeholders[offsetC + 2*k +j].text = " "
                slide.placeholders[offsetD + 2*k +j].text = " "
        k += 1

    if k<7:
        for n in range(k, 7):
            slide.placeholders[offsetA + n].text = " "
            for i in range(0, 2):
                slide.placeholders[offsetB + 2*n +i].text = " "
                slide.placeholders[offsetC + 2*n +i].text = " "
                slide.placeholders[offsetD + 2*n +i].text = " "
    return prs


# Adds slides for the "Fill in the Blanks" Exercise 
# First slide is up to 30 vocabulary to choose from
# Following slides are sentences with blanks in them to be filled in
def addslide_FITB(prs, data):

    FITB_questions = data["FITB_questions"]
    random.shuffle(FITB_questions)

    # GENERATE VOCAB SLIDE
    slide_layout = prs.slide_layouts[3]
    slide = prs.slides.add_slide(slide_layout)

    for i in range(0, 30):
        try:
            sentence = FITB_questions[i]
            slide.placeholders[13 + i].text = "□" + sentence.get("vocab")
        except:
            slide.placeholders[13 + i].text = " "

    # GENERATE SENTENCE SLIDES
    for i, eachline in enumerate(FITB_questions[:30]):
        if i%7 == 0:
            slide_layout = prs.slide_layouts[4 + i%3]
            slide = prs.slides.add_slide(slide_layout)
            k = 0
        try:
            slide.placeholders[13 + k].text = "✎"+eachline.get("sentence").replace("______", "_________")
        except:
            slide.placeholders[13 + k].text = " "
        k += 1
    
    if k<7:
        for n in range(k, 7):
            slide.placeholders[13 + n].text = " "


# Following slides are sentences with blanks in them to be filled in
def addslide_twoimages(prs, image1, image2):
    slide_layout = prs.slide_layouts[7]

    # Place image1
    slide = prs.slides.add_slide(slide_layout)
    img = Image.open(image1)
    aspect_ratio = img.width / img.height
    left = prs.slide_width/2
    top = Inches(0.1)
    height = prs.slide_height - top * 2
    width = height * aspect_ratio
    slide.shapes.add_picture(image1, left, top, width, height)

    # Place image2
    try:
        img = Image.open(image2)
        aspect_ratio = img.width / img.height
        left = prs.slide_width/2 
        top = Inches(0.1)
        height = prs.slide_height - top * 2
        width = height * aspect_ratio
        left = prs.slide_width/2 - width
        slide.shapes.add_picture(image2, left, top, width, height)
    except:
        print("completing comic image placements with only 1 image on final page")

    
# These slides are for the chinese-english meaning matching exercise
def addslide_englishmatch(prs, data):
    # GENERATE VOCAB SLIDE
    slide_layout = prs.slide_layouts[8]
    max_placeholders = 14
    slides_needed = len(data) // max_placeholders + 1
    vocab_str_list = list(data)

    for slide_num in range(slides_needed):
        slide = prs.slides.add_slide(slide_layout)
        english_array = []

        for i in range(0, max_placeholders):
            single_vocab = vocab_str_list[i + 2*slide_num]
            try:
                

                slide.placeholders[13 + i].text = "______" + single_vocab
                english_array.append(data[single_vocab]["english_meaning"])
            except:
                slide.placeholders[13 + i].text = " "
        
        random.shuffle(english_array)
        for j in range(0, 14):
            try:
                slide.placeholders[27 + j].text = "□ "+str(j+1) +":  "+ english_array[j]
            except:
                slide.placeholders[13 + j].text = " "

# These slides are for the sentence assembly exercise 
def addslide_unjumble(prs, data):
    # GENERATE VOCAB SLIDE
    slide_layout = prs.slide_layouts[9]
    max_placeholders = 5
    #print(data)
    FITB_questions = data["FITB_questions"]
    seg_data = [item for item in FITB_questions if item.get("seg") == True]
    slides_needed = len(seg_data) // max_placeholders + 1

    for slide_num in range(slides_needed):
        slide = prs.slides.add_slide(slide_layout)

        for i in range(0, max_placeholders):
            try:
                single_sentence = seg_data[i + 5*slide_num].get("complete_sentence")
                if "：" in single_sentence:
                    single_sentence = single_sentence.split("：", 1)[1]
                jumble_sentence = jieba.lcut(single_sentence)
                random.shuffle(jumble_sentence)
                #print(jumble_sentence)
                slide.placeholders[13 + i].text = " 🍡 ".join(jumble_sentence)
            except:
                slide.placeholders[13 + i].text = " "
        
# For unit testing
if __name__ == "__main__":
    template_path = "powerpoint_templates/doraemon_lesson_template_v1.pptx"
    output_path = "output_lesson/unittest.pptx"
    prs = Presentation(template_path)
    addslide_twoimages(prs, "doraemon_content/story2/images/1.png", "doraemon_content/story2/images/2.png")
    # Save the new presentation
    prs.save(output_path)
