import json
from pptx import Presentation
import random
import datetime

from utils_genslides import addslide_title_teacher, addslide_title_student, addslide_comicscript, addslide_FITB, addslide_twoimages, addslide_englishmatch, addslide_unjumble

# Load json file with the lesson content into a dictionary
def load_json(json_file):
    with open(json_file) as json_file:
        lesson_content = json.load(json_file)
    return lesson_content

def load_json_with_dict(json_file):
    with open(json_file) as file:
        try:
            data = json.load(file)
            if isinstance(data, dict):
                return data
            else:
                raise ValueError("JSON content is not a dictionary")
        except json.JSONDecodeError as e:
            print(f"Error decoding JSON: {e}")
            return None

def create_presentation(img_path, img_count, template_path, output_path, script_content, FITB_content, vocab_content):
    # Load the template presentation
    prs = Presentation(template_path)

    # Add the title slides
    try:
        print("Loading JSON files")
        data_script = load_json_with_dict(script_content)
        data_vocab = load_json_with_dict(vocab_content)
        data_FITB = load_json_with_dict(FITB_content)
    except:
        print("Error loading JSON files")
        print("I need three input files in the final_json folder: aggregated_script.json, FITB_final.json, and Vocab_final.json")
        print("please make sure the files are in the correct folder")
        print("exiting...")
        return

    # Add the title slides
    addslide_title_teacher(prs, data_script)
    addslide_title_student(prs, data_script)

    # Add the comic strip slides
    for i in range(1, ((img_count + 1) // 2) +1):
        image1 = img_path + str(i*2 -1 ) + ".png"
        image2 = img_path + str(i*2) + ".png"
        addslide_twoimages(prs, image1, image2)

    # Add the script & vocab slides
    addslide_comicscript(prs, data_script)

    # Add the vocabulary selection slide and the corresponding fill-in-the-blanks slides
    addslide_FITB(prs, data_FITB)

    # Add the chinese-english matching slides
    addslide_englishmatch(prs, data_vocab["vocab_list"])

    # Add the unjumble sentences slides
    addslide_unjumble(prs, data_FITB)

    # Save the new presentation
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")


# For new templates, use this to find out the correct reference index for each placeholder
def map_placeholder_index(template_path, output_path):
    # Load the template presentation
    prs = Presentation(template_path)
    for each_layout in range(0,9):
        slide_layout = prs.slide_layouts[each_layout]
        slide = prs.slides.add_slide(slide_layout)
        for shape in slide.placeholders:
            content = slide.placeholders[shape.placeholder_format.idx]
            content.text = "IDX: "+str(shape.placeholder_format.idx)
    # Save the new presentation
    prs.save(output_path)
    print(f"Index deck saved to: {output_path}")